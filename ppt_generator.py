import json
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE



# --- 1. Design Constants ---
# Colors
GOOGLE_BLUE = RGBColor(66, 133, 244)
TEXT_COLOR = RGBColor(32, 33, 36)
DARK_GRAY = RGBColor(64, 64, 64) # Added dark gray color
WHITE = RGBColor(255, 255, 255) # Added white color
RED = RGBColor(255, 0, 0)
ORANGE = RGBColor(255, 165, 0)
GREEN = RGBColor(0, 128, 0)
# Slide Dimensions (16:9)
SLIDE_WIDTH = Inches(16)
SLIDE_HEIGHT = Inches(9)
# Margins
MARGIN_LEFT = Inches(1.0)
MARGIN_RIGHT = Inches(1.0)
MARGIN_TOP = Inches(0.8)
MARGIN_BOTTOM = Inches(0.8)
# Font Sizes
FONT_HEADLINE = 'Meiryo UI'
FONT_BODY = 'Meiryo UI'
TITLE_FONT_SIZE = Pt(60)
SECTION_TITLE_FONT_SIZE = Pt(54) # This will be handled dynamically in draw_section_slide
SLIDE_TITLE_FONT_SIZE = Pt(36)
SUBHEAD_FONT_SIZE = Pt(30)
BODY_FONT_SIZE = Pt(20)
PROCESS_STEP_NUMBER_FONT_SIZE = Pt(12)
PROCESS_STEP_DESC_FONT_SIZE = Pt(10)

# --- 2. Helper Functions ---

def add_speaker_notes(slide, notes_text):
    """Adds speaker notes to the slide."""
    if notes_text:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        
        # Truncate notes_text to avoid python-pptx internal naming issues
        MAX_NOTES_LENGTH_FOR_NAME = 1000
        if len(notes_text) > MAX_NOTES_LENGTH_FOR_NAME:
            notes_text = notes_text[:MAX_NOTES_LENGTH_FOR_NAME] + "..." # Add ellipsis

        text_frame.text = notes_text



def set_formatted_text_in_frame(text_frame, text):
    """
    Clears a text_frame and populates it by parsing text with **bold** 
    and [[highlight]] syntax.
    """
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    # ガード節を追加。textがNoneや空の場合は、クリアされた状態で終了
    if not text:
        return
    # Truncate text to avoid python-pptx internal naming issues
    MAX_TEXT_LENGTH_FOR_NAME = 1000
    if len(text) > MAX_TEXT_LENGTH_FOR_NAME:
        text = text[:MAX_TEXT_LENGTH_FOR_NAME] + "..." # Add ellipsis for truncation

    apply_formatted_text_to_paragraph(p, text)

def apply_formatted_text_to_paragraph(p, text):
    """
    Parses text with **bold** and [[highlight]] syntax and adds it as runs
    to a paragraph object.
    """
    if not text:
        return
    # Split text by markers, keeping the markers
    parts = re.split(r'(\*\*.*?\*\*|\[\[.*?\]\])', text)

    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            run = p.add_run()
            run.text = part[2:-2]
            run.font.bold = True
            run.font.name = FONT_BODY
        elif part.startswith('[[') and part.endswith(']]'):
            run = p.add_run()
            run.text = part[2:-2]
            run.font.bold = True
            run.font.color.rgb = GOOGLE_BLUE
            run.font.name = FONT_BODY
        elif part:
            run = p.add_run()
            run.text = part
            run.font.name = FONT_BODY

# --- 3. Slide Drawing Functions ---

def draw_title_slide(slide, data):
    """Draws a title slide."""
    # Title
    title_shape = slide.shapes.add_textbox(
        Inches(1), Inches(3.5), SLIDE_WIDTH - Inches(2), Inches(2)
    )
    title_tf = title_shape.text_frame
    title_tf.word_wrap = True
    p = title_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = TITLE_FONT_SIZE
    p.font.name = FONT_HEADLINE
    p.font.color.rgb = TEXT_COLOR
    p.font.bold = True # Added bold
    # For the main title, we set the whole frame
    set_formatted_text_in_frame(title_tf, data.get("title", ""))
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER # Re-apply alignment

    # Date
    date_shape = slide.shapes.add_textbox(
        Inches(1), Inches(8), SLIDE_WIDTH - Inches(2), Inches(0.5)
    )
    date_tf = date_shape.text_frame
    p = date_tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(24)
    p.font.name = FONT_BODY
    p.text = data.get("date", "")
    print(f"  - Drawing Title Slide: {data.get('title','')}")

def draw_section_slide(slide, data):
    """Draws a section header slide with section number and title side-by-side."""
    section_no = data.get("sectionNo", "")
    section_title_text = data.get("title", "")

    section_title_text = re.sub(r'^\d+[\.\:\-]?\s*', '', section_title_text)
    # Calculate positions for two text boxes
    total_width = SLIDE_WIDTH - MARGIN_LEFT * 2
    # Estimate width for section number (rough estimate based on font size)
    # This might need fine-tuning based on actual number of digits
    section_no_width = Inches(2.0) # Adjust as needed
    section_title_width = total_width - section_no_width - Inches(0.5) # 0.5 inch gap

    # Section Number
    no_shape = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(3), section_no_width, Inches(3)
    )
    no_tf = no_shape.text_frame
    no_tf.word_wrap = True
    no_tf.vertical_anchor = MSO_ANCHOR.MIDDLE # Added vertical alignment
    p_no = no_tf.paragraphs[0]
    p_no.alignment = PP_ALIGN.RIGHT # Align number to right within its box
    p_no.font.size = Pt(150) # Requested size
    p_no.font.name = FONT_HEADLINE
    p_no.font.color.rgb = GOOGLE_BLUE # Requested color
    p_no.font.bold = True # Requested bold
    p_no.text = str(section_no)

    # Section Title
    title_shape = slide.shapes.add_textbox(
        MARGIN_LEFT + section_no_width + Inches(0.5), Inches(3), section_title_width, Inches(3)
    )
    title_tf = title_shape.text_frame
    title_tf.word_wrap = True
    title_tf.vertical_anchor = MSO_ANCHOR.MIDDLE # Added vertical alignment
    p_title = title_tf.paragraphs[0]
    p_title.alignment = PP_ALIGN.LEFT # Align title to left within its box
    p_title.font.size = Pt(60) # Requested size
    p_title.font.name = FONT_HEADLINE
    p_title.font.color.rgb = TEXT_COLOR
    p_title.font.bold = True # Requested bold
    set_formatted_text_in_frame(title_tf, section_title_text)

    print(f"  - Drawing Section Slide: {section_title_text}")

def draw_content_slide(slide, data):
    """Draws a content slide (1-column or 2-column)."""
    # Slide Title
    title_shape = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(1)
    )
    title_tf = title_shape.text_frame
    title_tf.word_wrap = True
    p = title_tf.paragraphs[0]
    p.font.size = SLIDE_TITLE_FONT_SIZE
    p.font.name = FONT_HEADLINE
    p.font.color.rgb = TEXT_COLOR
    p.font.bold = True
    set_formatted_text_in_frame(title_tf, data.get("title", ""))

    body_top = MARGIN_TOP + Inches(1.2)
    body_height = SLIDE_HEIGHT - body_top - MARGIN_BOTTOM

    if data.get("twoColumn"):
        # Two-column layout
        columns_data = data.get("columns", [[], []])
        left_column_items = columns_data[0] if len(columns_data) > 0 else []
        right_column_items = columns_data[1] if len(columns_data) > 1 else []

        column_width = (SLIDE_WIDTH - MARGIN_LEFT * 2 - Inches(0.5)) / 2 # 0.5 inch gap between columns

        # Left Column
        left_body_shape = slide.shapes.add_textbox(
            MARGIN_LEFT, body_top, column_width, body_height
        )
        left_body_tf = left_body_shape.text_frame
        left_body_tf.clear()
        left_body_tf.word_wrap = True
        for item_text in left_column_items:
            p = left_body_tf.add_paragraph()
            p.font.size = BODY_FONT_SIZE
            p.font.name = FONT_BODY
            p.level = 1
            apply_formatted_text_to_paragraph(p, item_text)

        # Right Column
        right_body_shape = slide.shapes.add_textbox(
            MARGIN_LEFT + column_width + Inches(0.5), body_top, column_width, body_height
        )
        right_body_tf = right_body_shape.text_frame
        right_body_tf.clear()
        right_body_tf.word_wrap = True
        for item_text in right_column_items:
            p = right_body_tf.add_paragraph()
            p.font.size = BODY_FONT_SIZE
            p.font.name = FONT_BODY
            p.level = 1
            apply_formatted_text_to_paragraph(p, item_text)
    else:
        # One-column layout (existing logic)
        body_shape = slide.shapes.add_textbox(
            MARGIN_LEFT, body_top, SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, body_height
        )
        body_tf = body_shape.text_frame
        body_tf.clear() # Clear default paragraph
        body_tf.word_wrap = True

        # Subhead
        if "subhead" in data:
            p = body_tf.add_paragraph()
            p.font.size = SUBHEAD_FONT_SIZE
            p.font.name = FONT_BODY
            apply_formatted_text_to_paragraph(p, data["subhead"])
            p.space_after = Pt(12)

        # Points (bullets)
        if "points" in data:
            slide_title_str = data.get("title", "")
            is_agenda_slide = "アジェンダ" in slide_title_str or "Agenda" in slide_title_str
            for point_text in data["points"]:
                p = body_tf.add_paragraph()
                p.font.size = Pt(36) # Changed to Pt(28)
                p.font.name = FONT_BODY
                p.level = 1

                if is_agenda_slide:
                    p.font.bold = True # アジェンダスライドなら段落全体を太字にする

                apply_formatted_text_to_paragraph(p, point_text)

                if is_agenda_slide:
                    p.font.bold = True # 念のためヘルパー関数の後にも設定
        
    print(f"  - Drawing Content Slide: {data.get('title','')}")

def draw_compare_slide(slide, data):
    """Draws a compare slide (two columns)."""
    # Slide Title
    title_shape = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(1)
    )
    title_tf = title_shape.text_frame
    title_tf.word_wrap = True
    p = title_tf.paragraphs[0]
    p.font.size = SLIDE_TITLE_FONT_SIZE
    p.font.name = FONT_HEADLINE
    p.font.color.rgb = TEXT_COLOR
    p.font.bold = True
    set_formatted_text_in_frame(title_tf, data.get("title", ""))

    # Subhead
    body_top = MARGIN_TOP + Inches(1.2)
    if "subhead" in data:
        subhead_shape = slide.shapes.add_textbox(
            MARGIN_LEFT, body_top, SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(0.5)
        )
        subhead_tf = subhead_shape.text_frame
        p = subhead_tf.paragraphs[0]
        p.font.size = Pt(24) # Changed to Pt(24)
        p.font.name = FONT_BODY
        p.font.color.rgb = GOOGLE_BLUE # Changed to GOOGLE_BLUE
        apply_formatted_text_to_paragraph(p, data["subhead"])
        body_top += Inches(0.6)

    # Column Layout
    column_width = (SLIDE_WIDTH - MARGIN_LEFT * 2 - Inches(0.5)) / 2 # 0.5 inch gap between columns
    column_height = SLIDE_HEIGHT - body_top - MARGIN_BOTTOM

    # Left Column Title
    left_title_shape = slide.shapes.add_textbox(
        MARGIN_LEFT, body_top, column_width, Inches(0.5)
    )
    left_title_shape.fill.solid() # Added solid fill
    left_title_shape.fill.fore_color.rgb = DARK_GRAY # Added dark gray fill
    left_title_tf = left_title_shape.text_frame
    p = left_title_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER # Added center alignment
    p.font.size = Pt(30) # Changed to Pt(30)
    p.font.name = FONT_HEADLINE
    p.font.color.rgb = WHITE # Changed to WHITE
    p.font.bold = True
    set_formatted_text_in_frame(left_title_tf, data.get("leftTitle", ""))

    # Right Column Title
    right_title_shape = slide.shapes.add_textbox(
        MARGIN_LEFT + column_width + Inches(0.5), body_top, column_width, Inches(0.5)
    )
    right_title_shape.fill.solid() # Added solid fill
    right_title_shape.fill.fore_color.rgb = DARK_GRAY # Added dark gray fill
    right_title_tf = right_title_shape.text_frame
    p = right_title_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER # Added center alignment
    p.font.size = Pt(30) # Changed to Pt(30)
    p.font.name = FONT_HEADLINE
    p.font.color.rgb = WHITE # Changed to WHITE
    p.font.bold = True
    set_formatted_text_in_frame(right_title_tf, data.get("rightTitle", ""))

    # Left Column Items
    left_items_top = body_top + Inches(0.6)
    left_items_shape = slide.shapes.add_textbox(
        MARGIN_LEFT, left_items_top, column_width, column_height - Inches(0.6)
    )
    left_items_tf = left_items_shape.text_frame
    left_items_tf.clear()
    for item in data.get("leftItems", []):
        p = left_items_tf.add_paragraph()
        p.font.size = Pt(24) # Changed to Pt(24)
        p.font.name = FONT_BODY
        p.level = 1
        apply_formatted_text_to_paragraph(p, item)

    # Right Column Items
    right_items_top = body_top + Inches(0.6)
    right_items_shape = slide.shapes.add_textbox(
        MARGIN_LEFT + column_width + Inches(0.5), right_items_top, column_width, column_height - Inches(0.6)
    )
    right_items_tf = right_items_shape.text_frame
    right_items_tf.clear()
    for item in data.get("rightItems", []):
        p = right_items_tf.add_paragraph()
        p.font.size = Pt(24) # Changed to Pt(24)
        p.font.name = FONT_BODY
        p.level = 1
        apply_formatted_text_to_paragraph(p, item)

    print(f"  - Drawing Compare Slide: {data.get('title','')}")



def draw_process_slide(slide, data):
    """Draws a process slide (horizontal flow of steps)."""
    # Slide Title
    title_shape = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(1)
    )
    title_tf = title_shape.text_frame
    title_tf.word_wrap = True
    p = title_tf.paragraphs[0]
    p.font.size = SLIDE_TITLE_FONT_SIZE
    p.font.name = FONT_HEADLINE
    p.font.color.rgb = TEXT_COLOR
    p.font.bold = True
    set_formatted_text_in_frame(title_tf, data.get("title", ""))

    # Subhead
    body_top = MARGIN_TOP + Inches(1.2)
    if "subhead" in data:
        subhead_shape = slide.shapes.add_textbox(
            MARGIN_LEFT, body_top, SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(0.5)
        )
        subhead_tf = subhead_shape.text_frame
        p = subhead_tf.paragraphs[0]
        p.font.size = Pt(24) # Changed to Pt(24)
        p.font.name = FONT_BODY
        p.font.color.rgb = GOOGLE_BLUE # Changed to GOOGLE_BLUE
        apply_formatted_text_to_paragraph(p, data["subhead"])
        body_top += Inches(0.6)

    # Process Steps
    steps = data.get("steps", [])
    if not steps: return

    num_steps = len(steps)
    
    # Define fixed step size and connector width
    step_size = Inches(1.5) # Make steps square
    connector_width = Inches(0.8) # Width of the connector line

    # Calculate total width available for steps and gaps
    available_width = SLIDE_WIDTH - MARGIN_LEFT * 2

    # Calculate total width occupied by steps and connectors
    total_steps_width = num_steps * step_size
    total_connectors_width = (num_steps - 1) * connector_width

    # Calculate remaining width for gaps
    remaining_width_for_gaps = available_width - total_steps_width - total_connectors_width

    # Calculate individual gap width
    if num_steps > 1:
        total_num_gaps = (num_steps - 1) * 2 # Two gaps per connector
        gap_width = remaining_width_for_gaps / total_num_gaps
    else:
        gap_width = Inches(0) # No gaps if only one step

    # Calculate vertical centering for the process steps
    diagram_height = step_size # The height of the process flow is essentially the height of a step
    available_vertical_space = SLIDE_HEIGHT - body_top - MARGIN_BOTTOM
    start_y = body_top + (available_vertical_space - diagram_height) / 2

    current_x = MARGIN_LEFT

    for i, step_text in enumerate(steps):
        # Draw Step Box
        step_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, current_x, start_y, step_size, step_size) # Changed to RECTANGLE and step_size
        step_shape.fill.solid()
        step_shape.fill.fore_color.rgb = DARK_GRAY # Dark gray fill
        step_tf = step_shape.text_frame
        step_tf.word_wrap = True
        step_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = step_tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(24) # Changed to Pt(24)
        p.font.name = FONT_BODY
        apply_formatted_text_to_paragraph(p, step_text)

        current_x += step_size

        # Draw Arrow (if not last step)
        if i < num_steps - 1:
            arrow_start_x = current_x + gap_width
            arrow_end_x = arrow_start_x + connector_width
            arrow_y = start_y + step_size / 2
            
            # 矢印図形を追加
            arrow_shape = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, # 右矢印の形状
                arrow_start_x, arrow_y - Inches(0.1), # 矢印の高さの半分を引いて中央に配置
                connector_width, Inches(0.2) # 矢印の幅と高さを設定
            )
            arrow_shape.fill.solid()
            arrow_shape.fill.fore_color.rgb = TEXT_COLOR # 矢印の色を設定
            arrow_shape.line.fill.background() # 線の色を塗りつぶしと同じにする（線なし）

            current_x += gap_width + connector_width + gap_width

    print(f"  - Drawing Process Slide: {data.get('title','')}")



def draw_closing_slide(slide, data):
    """Draws a simple closing slide with 'Thank you'."""
    # Title (Thank you)
    title_shape = slide.shapes.add_textbox(
        Inches(1), Inches(3.5), SLIDE_WIDTH - Inches(2), Inches(2)
    )
    title_tf = title_shape.text_frame
    title_tf.word_wrap = True
    p = title_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = TITLE_FONT_SIZE # Use the same font size as main title
    p.font.name = FONT_HEADLINE
    p.font.color.rgb = TEXT_COLOR
    p.font.bold = True
    p.text = "Thank you" # Directly set text as it's fixed

    print(f"  - Drawing Closing Slide: Thank you")

def draw_table_slide(slide, data):
    """Draws a table slide."""
    # Slide Title
    title_shape = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(1)
    )
    title_tf = title_shape.text_frame
    title_tf.word_wrap = True
    p = title_tf.paragraphs[0]
    p.font.size = SLIDE_TITLE_FONT_SIZE
    p.font.name = FONT_HEADLINE
    p.font.color.rgb = TEXT_COLOR
    p.font.bold = True
    set_formatted_text_in_frame(title_tf, data.get("title", ""))

    # Subhead
    body_top = MARGIN_TOP + Inches(1.2)
    if "subhead" in data:
        subhead_shape = slide.shapes.add_textbox(
            MARGIN_LEFT, body_top, SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(0.5)
        )
        subhead_tf = subhead_shape.text_frame
        p = subhead_tf.paragraphs[0]
        p.font.size = Pt(24) # Changed to Pt(24)
        p.font.name = FONT_BODY
        p.font.color.rgb = GOOGLE_BLUE # Changed to GOOGLE_BLUE
        apply_formatted_text_to_paragraph(p, data["subhead"])
        body_top += Inches(0.6)

    # Table
    headers = data.get("headers", [])
    rows_data = data.get("rows", [])

    if not headers and not rows_data:
        print("    - WARNING: No headers or rows data for table slide. Skipping table drawing.")
        return

    num_cols = len(headers) if headers else (len(rows_data[0]) if rows_data else 0)
    num_rows = len(rows_data) + (1 if headers else 0) # +1 for header row

    if num_cols == 0 or num_rows == 0:
        print("    - WARNING: Invalid table dimensions. Skipping table drawing.")
        return

    table_left = MARGIN_LEFT
    table_top = body_top + Inches(0.2) # Small gap after subhead
    table_width = SLIDE_WIDTH - MARGIN_LEFT * 2
    table_height = SLIDE_HEIGHT - table_top - MARGIN_BOTTOM

    # Add table
    table_shape = slide.shapes.add_table(num_rows, num_cols, table_left, table_top, table_width, table_height)
    table = table_shape.table

    # Set column widths (equal distribution for now)
    for col_idx in range(num_cols):
        table.columns[col_idx].width = Emu(table_width / num_cols)

    # Populate headers
    if headers:
        for col_idx, header_text in enumerate(headers):
            cell = table.cell(0, col_idx)
            text_frame = cell.text_frame
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE # Set vertical alignment to middle
            p = text_frame.paragraphs[0]
            apply_formatted_text_to_paragraph(p, header_text) 
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            p.font.size = BODY_FONT_SIZE # Use body font size for headers

    # Populate rows
    start_row_idx = 1 if headers else 0
    for row_idx, row_data in enumerate(rows_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(start_row_idx + row_idx, col_idx)
            text_frame = cell.text_frame
            p = text_frame.paragraphs[0]
            apply_formatted_text_to_paragraph(p, cell_text)
            p.alignment = PP_ALIGN.LEFT # Default to left alignment for body cells
            p.font.size = BODY_FONT_SIZE

    print(f"  - Drawing Table Slide: {data.get('title','')}")



def draw_timeline_slide(slide, data):
    """
    Draws a DYNAMIC timeline slide, alternating milestones above and below the line
    based on the 'milestones' array in the data.
    """
    # --- 1. Title and Subhead (Standard) ---
    title_text = data.get("title", "Timeline")
    title_shape = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(1)
    )
    title_tf = title_shape.text_frame
    p = title_tf.paragraphs[0]
    p.font.size = SLIDE_TITLE_FONT_SIZE
    p.font.name = FONT_HEADLINE
    p.font.color.rgb = TEXT_COLOR
    p.font.bold = True
    set_formatted_text_in_frame(title_tf, title_text)

    body_top = MARGIN_TOP + Inches(1.2)
    if "subhead" in data:
        desc_shape = slide.shapes.add_textbox(
            MARGIN_LEFT, body_top, SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(0.8)
        )
        desc_tf = desc_shape.text_frame
        p = desc_tf.paragraphs[0]
        p.font.size = BODY_FONT_SIZE
        p.font.name = FONT_BODY
        p.font.color.rgb = TEXT_COLOR
        set_formatted_text_in_frame(desc_tf, data["subhead"])
        body_top += Inches(1.0) # Add space for the description
    else:
        body_top += Inches(0.2) # Smaller gap if no subhead

    # --- 2. Timeline Horizontal Bar ---
    # Center the timeline vertically in the remaining space
    available_height = SLIDE_HEIGHT - body_top - MARGIN_BOTTOM
    timeline_y = body_top + (available_height / 2) # Vertical center

    timeline_width = SLIDE_WIDTH - MARGIN_LEFT * 2
    timeline_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN_LEFT, timeline_y - Inches(0.025), # Center the line thickness
        timeline_width, Inches(0.05) # Thickness of the line
    )
    timeline_shape.fill.solid()
    timeline_shape.fill.fore_color.rgb = TEXT_COLOR # Black
    timeline_shape.line.fill.background() # No border

    # --- 3. Dynamic Milestones Processing ---
    milestones = data.get("milestones", [])
    num_milestones = len(milestones)
    if num_milestones == 0:
        print(f"  - WARNING: No milestones data provided for timeline slide.")
        return

    # Calculate horizontal spacing
    # We divide the width into (N+1) segments to get N points evenly spaced
    timeline_usable_width = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
    if num_milestones > 1:
        spacing = timeline_usable_width / (num_milestones - 1) # Space between N points
        start_x = MARGIN_LEFT
    else:
        spacing = 0 # Only one item
        start_x = MARGIN_LEFT + (timeline_usable_width / 2) # Center if only one


    marker_size = Inches(0.5)
    desc_box_width = Inches(3.5) # Fixed width for description wrapping
    desc_box_height = Inches(1.5) # Fixed height

    # Define the Y positions (fixed heights for "above" and "below")
    date_y_above = timeline_y - Inches(0.7) # Date position (relative to timeline)
    desc_y_below = timeline_y + Inches(0.7) # Description start Y (below line)
    
    # Calculate Y pos for "above": Top of box must be (Timeline Y) - (Gap) - (Box Height)
    desc_y_above = timeline_y - Inches(0.7) - desc_box_height 


    for i, milestone in enumerate(milestones):
        milestone_label = milestone.get("label", "")
        milestone_date = milestone.get("date", "")
        
        # Calculate marker X position
        if num_milestones > 1:
            marker_center_x = start_x + (i * spacing)
        else:
            marker_center_x = start_x # Use the centered start_x

        # --- Draw Marker Circle ---
        marker_x = marker_center_x - (marker_size / 2)
        marker_y = timeline_y - (marker_size / 2)
        marker_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, marker_x, marker_y, marker_size, marker_size
        )
        marker_shape.fill.solid()
        marker_shape.fill.fore_color.rgb = GOOGLE_BLUE # Blue
        marker_shape.line.fill.background() # No border

        # --- Calculate description/date positions (ALTERNATING LOGIC) ---
        
        current_date_y = date_y_above # Dates always go above the line
        
        if i % 2 == 0:
            # Even index (0, 2, ...): Place description BELOW
            current_desc_y = desc_y_below
        else:
            # Odd index (1, 3, ...): Place description ABOVE
            current_desc_y = desc_y_above
            # Since description is above, move date above the description box
            current_date_y = current_desc_y - Inches(0.4) # Place date above the upper description


        # --- Draw Date Label ---
        date_shape = slide.shapes.add_textbox(
            marker_center_x - (desc_box_width / 2), current_date_y, desc_box_width, Inches(0.5)
        )
        date_tf = date_shape.text_frame
        p_date = date_tf.paragraphs[0]
        p_date.alignment = PP_ALIGN.CENTER
        p_date.font.size = Pt(24)
        p_date.font.name = FONT_BODY
        p_date.font.bold = True
        p_date.text = milestone_date

        # --- Draw Phase Description Text ---
        desc_x = marker_center_x - (desc_box_width / 2)
        
        desc_shape = slide.shapes.add_textbox(
            desc_x, current_desc_y, desc_box_width, desc_box_height
        )
        desc_tf = desc_shape.text_frame
        p_desc = desc_tf.paragraphs[0]
        p_desc.alignment = PP_ALIGN.CENTER
        p_desc.font.size = BODY_FONT_SIZE
        p_desc.font.name = FONT_BODY
        p_desc.text = milestone_label

    print(f"  - Drawing Dynamic Timeline Slide: {title_text}")

def draw_diagram_slide(slide, data):
    """Draws a diagram (lane) slide."""
    # Slide Title
    title_shape = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(1)
    )
    title_tf = title_shape.text_frame
    title_tf.word_wrap = True
    p = title_tf.paragraphs[0]
    p.font.size = SLIDE_TITLE_FONT_SIZE
    p.font.name = FONT_HEADLINE
    p.font.color.rgb = TEXT_COLOR
    p.font.bold = True
    set_formatted_text_in_frame(title_tf, data.get("title", ""))

    # Subhead
    body_top = MARGIN_TOP + Inches(1.2)
    if "subhead" in data:
        subhead_shape = slide.shapes.add_textbox(
            MARGIN_LEFT, body_top, SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(0.5)
        )
        subhead_tf = subhead_shape.text_frame
        p = subhead_tf.paragraphs[0]
        p.font.size = Pt(24) # Changed to Pt(24)
        p.font.name = FONT_BODY
        p.font.color.rgb = GOOGLE_BLUE # Changed to GOOGLE_BLUE
        apply_formatted_text_to_paragraph(p, data["subhead"])
        body_top += Inches(0.6)

    # Lanes
    lanes = data.get("lanes", [])
    if not lanes: return

    num_lanes = len(lanes)
    lane_height = (SLIDE_HEIGHT - body_top - MARGIN_BOTTOM) / num_lanes
    lane_width = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT

    current_y = body_top

    for i, lane_data in enumerate(lanes):
        lane_title = lane_data.get("title", "")
        lane_items = lane_data.get("items", [])

        # Lane Title
        title_box = slide.shapes.add_textbox(
            MARGIN_LEFT, current_y, Inches(2), lane_height
        )
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = title_tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = SUBHEAD_FONT_SIZE
        p.font.name = FONT_HEADLINE
        p.font.bold = True
        p.text = lane_title

        # Lane Items
        items_box = slide.shapes.add_textbox(
            MARGIN_LEFT + Inches(2.2), current_y, lane_width - Inches(2.2), lane_height
        )
        items_tf = items_box.text_frame
        items_tf.clear()
        items_tf.word_wrap = True
        for item_text in lane_items:
            p = items_tf.add_paragraph()
            p.font.size = BODY_FONT_SIZE
            p.font.name = FONT_BODY
            p.level = 1
            apply_formatted_text_to_paragraph(p, item_text)

        current_y += lane_height

    print(f"  - Drawing Diagram Slide: {data.get('title','')}")

def draw_cards_slide(slide, data):
    """Draws a cards (grid) slide."""
    # Slide Title
    title_shape = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(1)
    )
    title_tf = title_shape.text_frame
    title_tf.word_wrap = True
    p = title_tf.paragraphs[0]
    p.font.size = SLIDE_TITLE_FONT_SIZE
    p.font.name = FONT_HEADLINE
    p.font.color.rgb = TEXT_COLOR
    p.font.bold = True
    set_formatted_text_in_frame(title_tf, data.get("title", ""))

    # Subhead
    body_top = MARGIN_TOP + Inches(1.2)
    if "subhead" in data:
        subhead_shape = slide.shapes.add_textbox(
            MARGIN_LEFT, body_top, SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(0.5)
        )
        subhead_tf = subhead_shape.text_frame
        p = subhead_tf.paragraphs[0]
        p.font.size = Pt(24) # Changed to Pt(24)
        p.font.name = FONT_BODY
        p.font.color.rgb = GOOGLE_BLUE # Changed to GOOGLE_BLUE
        apply_formatted_text_to_paragraph(p, data["subhead"])
        body_top += Inches(0.6)

    # Cards
    items = data.get("items", [])
    if not items: return

    num_columns = data.get("columns", 2) # Default to 2 columns
    if num_columns not in [2, 3]:
        print(f"    - WARNING: Invalid number of columns for cards slide: {num_columns}. Defaulting to 2.")
        num_columns = 2

    # Define card colors
    CARD_COLORS = [GOOGLE_BLUE, RED, ORANGE, GOOGLE_BLUE, GREEN, RED]

    # Calculate card dimensions
    total_width = SLIDE_WIDTH - MARGIN_LEFT * 2
    gap_between_cards = Inches(0.3)
    card_width = (total_width - (num_columns - 1) * gap_between_cards) / num_columns
    card_height = Inches(2.5) # Fixed height for cards

    start_y = body_top + Inches(0.5)
    current_x = MARGIN_LEFT
    current_y = start_y

    for i, item_data in enumerate(items):
        row_idx = i // num_columns
        col_idx = i % num_columns

        # Calculate position for the current card
        card_left = MARGIN_LEFT + col_idx * (card_width + gap_between_cards)
        card_top = start_y + row_idx * (card_height + gap_between_cards)

        # Card Shape
        card_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_left, card_top, card_width, card_height)
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = CARD_COLORS[i % len(CARD_COLORS)] # Cycle through defined colors

        # Card Content
        card_tf = card_shape.text_frame
        card_tf.word_wrap = True
        card_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        card_tf.clear()

        if isinstance(item_data, str):
            # Simple string item
            p = card_tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(24) # Changed to Pt(24)
            p.font.name = FONT_BODY
            p.font.bold = True # Added bold
            p.font.color.rgb = WHITE
            apply_formatted_text_to_paragraph(p, item_data)
        else:
            # Object with title and desc
            title = item_data.get("title", "")
            desc = item_data.get("desc", "")

            p_title = card_tf.add_paragraph()
            p_title.alignment = PP_ALIGN.CENTER
            p_title.font.size = Pt(24) # Changed to Pt(24)
            p_title.font.name = FONT_HEADLINE
            p_title.font.bold = True
            apply_formatted_text_to_paragraph(p_title, title)

            if desc:
                p_desc = card_tf.add_paragraph()
                p_desc.alignment = PP_ALIGN.CENTER
                p_desc.font.size = Pt(24) # Changed to Pt(24)
                p_desc.font.name = FONT_BODY
                p_desc.font.bold = True # Added bold
                apply_formatted_text_to_paragraph(p_desc, desc)

    print(f"  - Drawing Cards Slide: {data.get('title','')}")

def draw_progress_slide(slide, data):
    """Draws a progress slide."""
    # Slide Title
    title_shape = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(1)
    )
    title_tf = title_shape.text_frame
    title_tf.word_wrap = True
    p = title_tf.paragraphs[0]
    p.font.size = SLIDE_TITLE_FONT_SIZE
    p.font.name = FONT_HEADLINE
    p.font.color.rgb = TEXT_COLOR
    p.font.bold = True
    set_formatted_text_in_frame(title_tf, data.get("title", ""))

    # Subhead
    body_top = MARGIN_TOP + Inches(1.2)
    if "subhead" in data:
        subhead_shape = slide.shapes.add_textbox(
            MARGIN_LEFT, body_top, SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, Inches(0.5)
        )
        subhead_tf = subhead_shape.text_frame
        p = subhead_tf.paragraphs[0]
        p.font.size = Pt(24) # Changed to Pt(24)
        p.font.name = FONT_BODY
        p.font.color.rgb = GOOGLE_BLUE # Changed to GOOGLE_BLUE
        apply_formatted_text_to_paragraph(p, data["subhead"])
        body_top += Inches(0.6)

    # Progress Items
    items = data.get("items", [])
    if not items: return

    item_height = Inches(0.8)
    progress_bar_height = Inches(0.2)
    progress_bar_width = Inches(8) # Fixed width for progress bar

    current_y = body_top + Inches(0.5)

    for item in items:
        label = item.get("label", "")
        percent = item.get("percent", 0)

        # Label
        label_shape = slide.shapes.add_textbox(
            MARGIN_LEFT, current_y, Inches(4), item_height
        )
        label_tf = label_shape.text_frame
        p = label_tf.paragraphs[0]
        p.font.size = BODY_FONT_SIZE
        p.font.name = FONT_BODY
        p.font.bold = True # Added bold
        p.text = label

        # Progress Bar Background
        bar_bg_left = MARGIN_LEFT + Inches(4.5)
        bar_bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            bar_bg_left, current_y + (item_height - progress_bar_height) / 2,
            progress_bar_width, progress_bar_height
        )
        bar_bg_shape.fill.solid()
        bar_bg_shape.fill.fore_color.rgb = RGBColor(220, 220, 220) # Light gray
        bar_bg_shape.line.fill.background() # No border

        # Progress Bar Fill
        fill_width = progress_bar_width * (percent / 100.0)
        bar_fill_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            bar_bg_left, current_y + (item_height - progress_bar_height) / 2,
            fill_width, progress_bar_height
        )
        bar_fill_shape.fill.solid()
        bar_fill_shape.fill.fore_color.rgb = GOOGLE_BLUE # Google Blue
        bar_fill_shape.line.fill.background() # No border

        # Percentage Text
        percent_text_shape = slide.shapes.add_textbox(
            bar_bg_left + progress_bar_width + Inches(0.2), current_y, Inches(1), item_height
        )
        percent_text_tf = percent_text_shape.text_frame
        p = percent_text_tf.paragraphs[0]
        p.font.size = BODY_FONT_SIZE
        p.font.name = FONT_BODY
        p.text = f"{percent}%"

        current_y += item_height + Inches(0.2) # Add some spacing

    print(f"  - Drawing Progress Slide: {data.get('title','')}")


# --- 4. Main Execution Logic ---

def create_presentation(slides_data):
    """Creates a new presentation from slide data."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # Map types to functions
    slide_draw_functions = {
        "title": draw_title_slide,
        "section": draw_section_slide,
        "content": draw_content_slide,
        "compare": draw_compare_slide,
        "process": draw_process_slide,
        "closing": draw_closing_slide, # Added closing slide
        "table": draw_table_slide, # Added table slide
        "timeline": draw_timeline_slide, # Added timeline slide
        "diagram": draw_diagram_slide, # Added diagram slide
        "cards": draw_cards_slide, # Added cards slide
        "progress": draw_progress_slide, # Added progress slide
        # Other types will be added here
    }

    print("Starting presentation generation...")
    for i, slide_data in enumerate(slides_data):
        slide_type = slide_data.get("type")
        
        # --- NEW: Truncate title and notes for internal naming safety ---
        if "title" in slide_data and len(slide_data["title"]) > 1000:
            slide_data["title"] = slide_data["title"][:1000] + "..."
        if "notes" in slide_data and len(slide_data["notes"]) > 1000:
            slide_data["notes"] = slide_data["notes"][:1000] + "..."
        # --- END NEW ---
        print(f"Processing slide {i+1}: type='{slide_type}'")
        
        # Use a blank layout
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # --- NEW: Set slide.name explicitly ---
        slide.name = f"Slide_{i+1}_{slide_type}"
        # --- END NEW ---

        # Call the appropriate drawing function
        if slide_type in slide_draw_functions:
            slide_draw_functions[slide_type](slide, slide_data)
        else:
            print(f"  - WARNING: No drawing function for type '{slide_type}'. Skipping.")

        # Add speaker notes
        add_speaker_notes(slide, slide_data.get("notes"))

    return prs

